"""
PowerPoint (.pptx) loader.

Each slide becomes its own "page" (page_number = 1-based slide
number), matching how PDF pages work — a retrieved chunk can point
back to the exact slide it came from.
"""
from pptx import Presentation


def load_pptx(file_path):
    presentation = Presentation(file_path)
    pages = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
        pages.append((slide_index, "\n".join(texts)))

    total_chars = sum(len(t) for _, t in pages)
    print(f"✅ Extracted {total_chars} characters from {len(pages)} slide(s)")
    return pages